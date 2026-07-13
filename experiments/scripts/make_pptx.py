#!/usr/bin/env python3
"""Build the page-4 styled PPTX: task banner on top, two agent FPV videos below.

Layout replicates the reference deck slide 4 (960x540pt):
- 10pt navy (#1C2D54) strip on the left edge
- amber (#F29727) bold caption at (84, 40), 13pt
- navy (#182848) bold title at (84, 66), 24pt (the task given to agents)
- central visual area ~600x400 at y=108 -> two videos side by side
- gray page number bottom right
"""
import argparse
import json
import os

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x1C, 0x2D, 0x54)
NAVY_TEXT = RGBColor(0x18, 0x28, 0x48)
AMBER = RGBColor(0xF2, 0x97, 0x27)
GRAY = RGBColor(0xAA, 0xB0, 0xBA)
DARKGRAY = RGBColor(0x64, 0x6A, 0x73)
FONT = '思源黑体'

PT_PER_PX = 1.0  # reference deck uses 960x540 units == pt


def px(v):
    return Pt(v * PT_PER_PX)


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = FONT
    return box


def add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--clips-dir', required=True)
    ap.add_argument('--curve', help='score curve png (optional 2nd slide)')
    ap.add_argument('--case-video', help='MineCollab 10.1 original video mp4 (adds case slide first)')
    ap.add_argument('--case-poster', help='poster frame for the case video')
    ap.add_argument('--case-chat', help='telegram-style chat panel png')
    ap.add_argument('--out', required=True)
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width = px(960)
    prs.slide_height = px(540)
    blank = prs.slide_layouts[6]

    if args.case_video and args.case_chat:
        s0 = prs.slides.add_slide(blank)
        add_rect(s0, 0, 0, 10, 540, NAVY)
        add_text(s0, 84, 36, 780, 24, '背景 · 真实案例 | 对等协作中的执行现场（MineCollab 附录 10.1，arXiv 2504.17950）', 13, AMBER, bold=True)
        add_text(s0, 84, 62, 830, 34, '任务分好了，人人都在「好好干活」——石砖层却被队友拆光', 20, NAVY_TEXT, bold=True)
        s0.shapes.add_movie(args.case_video, px(60), px(118), px(330), px(330),
                            poster_frame_image=args.case_poster, mime_type='video/mp4')
        add_text(s0, 60, 452, 340, 44,
                 '▶ 论文原始录像（17s）：randy 成片挖掉 mandy 刚铺完的石砖层。', 10, DARKGRAY)
        s0.shapes.add_picture(args.case_chat, px(440), px(104), px(451), px(398))
        add_text(s0, 84, 512, 720, 20,
                 '注：MineCollab 的 agent 为对等自组织协作（无中央指挥）。互拆发生时人人自报「正常干活」——只要监督只看分工与最终验收，无论对等还是中央管理，这类过程损失都要到验收才暴露。', 10.5, DARKGRAY)

    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 10, 540, NAVY)
    add_text(s, 84, 40, 700, 26, '核心场景 · construction 互拆 · MINDcraft 双智能体实录', 13, AMBER, bold=True)
    add_text(s, 84, 66, 810, 34,
             '「Make a structure with the blueprint — share materials and build together」',
             19, NAVY_TEXT, bold=True)
    add_text(s, 84, 100, 800, 20,
             '任务下达后 andy / bob 自行协商分工；以下为冲突窗口内两人的第一视角（同一时刻，同一建筑）',
             11, DARKGRAY)

    # two videos side by side in the central visual area
    vw, vh = 380, 285  # 4:3 of 854x480 -> use 16:9: 380x214; pick 16:9
    vw, vh = 400, 225
    gap = 24
    total_w = vw * 2 + gap
    x0 = (960 - total_w) / 2 + 5
    y0 = 140
    for i, name in enumerate(['andy', 'bob']):
        clip = os.path.join(args.clips_dir, f'{name}_fpv.mp4')
        poster = os.path.join(args.clips_dir, f'{name}_poster.jpg')
        x = x0 + i * (vw + gap)
        s.shapes.add_movie(clip, px(x), px(y0), px(vw), px(vh),
                           poster_frame_image=poster, mime_type='video/mp4')
        add_rect(s, x, y0 + vh + 10, 8, 14, AMBER)
        add_text(s, x + 14, y0 + vh + 8, vw - 14, 18,
                 f'{name} · 第一视角', 12, NAVY_TEXT, bold=True)
    # annotate what to watch for
    info_path = os.path.join(args.clips_dir, 'clips_info.json')
    note = ''
    if os.path.exists(info_path):
        info = json.load(open(info_path))
        n = len(info.get('drops', []))
        sp = info.get('speed', 1.0)
        sptxt = f'，{sp:.1f}× 速' if sp > 1.001 else ''
        note = f'蓝图完成度在此窗口内 {n} 次下降：一方按过期蓝图快照「纠错」，拆除另一方刚铺好的方块，随后又被重建{sptxt}'
    add_text(s, 84, y0 + vh + 44, 800, 20, note, 11, DARKGRAY)
    add_text(s, 890, 508, 50, 20, '04', 11, GRAY)

    if args.curve and os.path.exists(args.curve):
        s2 = prs.slides.add_slide(blank)
        add_rect(s2, 0, 0, 10, 540, NAVY)
        add_text(s2, 84, 40, 700, 26, '证据 · 任务得分曲线（蓝图匹配度）', 13, AMBER, bold=True)
        add_text(s2, 84, 66, 800, 40, '「为什么每个人都干了很多活，我却还是看不到任务进展？」', 20, NAVY_TEXT, bold=True)
        s2.shapes.add_picture(args.curve, px(130), px(120), px(700), px(360))
        add_text(s2, 890, 508, 50, 20, '05', 11, GRAY)

    prs.save(args.out)
    print('saved', args.out)


if __name__ == '__main__':
    main()
